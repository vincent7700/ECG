# Genere un PowerPoint structure avec les metriques + exemples visuels
# Destine au partage avec l'equipe (Aymeric, Camille, Samuel, Ahmad, etc.)

import os, json, glob
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = r"c:\Users\v\Desktop\ECGPerturb-main"
METRIQUES_DIR = os.path.join(ROOT, "stage", "métriques")
OUTPUT_PPTX = os.path.join(METRIQUES_DIR, "metriques_modeles_ECG.pptx")

MODELS = [
    {
        "folder": "métriques training_npz_copie",
        "title": "Modele 1 : NPZ Full-Resolution (1024x1024)",
        "subtitle": "Detection des intersections de grille (5mm) - image entiere",
        "run_dir": r"data\training\runs_npz\run_20260518_164651",
        "best_dice": 0.7888,
        "best_iou": 0.6717,
        "best_epoch": 28,
    },
    {
        "folder": "métriques training_npz_patches_256",
        "title": "Modele 2 : NPZ Patch-Based (256x256)",
        "subtitle": "Detection des intersections de grille - approche par patches (PM Cardio)",
        "run_dir": r"data\training\runs_npz_patches256\run_20260520_170838",
        "best_dice": 0.6567,
        "best_iou": 0.5305,
        "best_epoch": 23,
    },
    {
        "folder": "métriques training_grid_major_png",
        "title": "Modele 3 : Grid Major (mask continu)",
        "subtitle": "Detection des lignes de grille majeures - mask binaire continu",
        "run_dir": r"data\training\runs\run_20260519_162837",
        "best_dice": 0.8491,
        "best_iou": 0.7532,
        "best_epoch": 29,
    },
]

FOLDERS = ["augmentation_brightness_160", "augmentation_rotation_15", "photos_crumbles"]
FOLDER_LABELS = {
    "augmentation_brightness_160": "Augmentation : Brightness +160%",
    "augmentation_rotation_15": "Augmentation : Rotation 15 deg",
    "photos_crumbles": "Photos : Papier froisse (crumbles)",
}

# Couleurs
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xD9, 0x53, 0x4F)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sw, sh = prs.slide_width, prs.slide_height

    # Bande haut
    band = slide.shapes.add_shape(1, 0, 0, sw, Inches(2.5))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    # Titre
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), sw - Inches(1), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Metriques des Modeles de Detection"
    r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), sw - Inches(1), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "Grille ECG - comparaison de 3 approches"
    r2.font.size = Pt(20); r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # Sous-titre bas
    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), sw - Inches(1), Inches(2))
    tf = tb3.text_frame
    tf.word_wrap = True

    lines = [
        ("Trois modeles compares :", True, 18, NAVY),
        ("    •  NPZ Full-Resolution 1024x1024", False, 14, GREY),
        ("    •  NPZ Patch-Based 256x256 (approche PM Cardio)", False, 14, GREY),
        ("    •  Grid Major (mask continu)", False, 14, GREY),
        ("", False, 14, GREY),
        ("Evaluation qualitative sur 3 sous-ensembles de PM Cardio (output_real) :", True, 16, NAVY),
        ("    •  augmentation_brightness_160", False, 14, GREY),
        ("    •  augmentation_rotation_15", False, 14, GREY),
        ("    •  photos_crumbles (papier froisse)", False, 14, GREY),
    ]
    for i, (text, bold, size, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color


def add_section_header(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    band = slide.shapes.add_shape(1, 0, Inches(2.5), sw, Inches(2.3))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), sw - Inches(1), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), sw - Inches(1), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(18); r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)


def add_slide_title(slide, text, sw):
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), sw - Inches(0.6), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = NAVY


def extract_hyperparams(resume_path):
    # Renvoie la liste des lignes de la section HYPERPARAMETRES du resume.txt.
    if not os.path.exists(resume_path):
        return []
    content = open(resume_path, encoding="utf-8").read()
    lines = []
    in_hp = False
    for ln in content.splitlines():
        if "HYPERPARAMETRES" in ln:
            in_hp = True
            continue
        if in_hp and ln.startswith("="):
            break
        if in_hp and ln.strip() and not ln.startswith("="):
            lines.append(ln.rstrip())
    return lines


def add_metrics_slide(prs, model):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    add_slide_title(slide, f"{model['title']} - Metriques", sw)

    # Cote gauche : metriques cles + hyperparametres
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(4.8), Inches(6))
    tf = tb.text_frame; tf.word_wrap = True

    head_lines = [
        ("RESULTATS", True, 16, NAVY),
        ("", False, 4, NAVY),
        (f"Best epoch        : {model['best_epoch']} / 30", False, 12, GREY),
        (f"Best Val Dice     : {model['best_dice']:.4f}", True, 12, ACCENT),
        (f"Best Val IoU      : {model['best_iou']:.4f}", True, 12, ACCENT),
        ("", False, 6, NAVY),
        ("HYPERPARAMETRES", True, 16, NAVY),
        ("", False, 4, NAVY),
    ]

    resume_path = os.path.join(METRIQUES_DIR, model["folder"], "resume.txt")
    hp_lines = extract_hyperparams(resume_path)

    all_lines = head_lines + [(ln, False, 11, GREY) for ln in hp_lines]

    for i, (text, bold, size, color) in enumerate(all_lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
        # Police monospace pour les hyperparametres (aligne les ":")
        if size == 11:
            run.font.name = "Consolas"

    # Cote droit : courbes d'entrainement
    curves_path = os.path.join(METRIQUES_DIR, model["folder"], "training_curves.png")
    if os.path.exists(curves_path):
        slide.shapes.add_picture(curves_path, Inches(5.4), Inches(1.5), width=Inches(7.6))

        cap = slide.shapes.add_textbox(Inches(5.4), Inches(5.7), Inches(7.6), Inches(0.4))
        p = cap.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "Courbes d'entrainement (loss + Dice/IoU)"
        r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = GREY


def _add_image_keep_ratio(slide, img_path, left, top, max_w, max_h):
    # Ajoute une image en respectant son ratio, centree dans la zone (max_w, max_h).
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        iw, ih = im.size
    ratio = iw / ih
    # Adapter aux limites
    target_w = max_w
    target_h = int(target_w / ratio)
    if target_h > max_h:
        target_h = max_h
        target_w = int(target_h * ratio)
    # Centrer
    cx = left + (max_w - target_w) // 2
    cy = top + (max_h - target_h) // 2
    slide.shapes.add_picture(img_path, cx, cy, width=target_w, height=target_h)


def add_examples_slides(prs, model, folder):
    # Cree 2 slides par dossier : slide 1 avec 2 images, slide 2 avec 1 image (plus grande).
    # On choisit 3 images sur les 5 disponibles : indices 0, 2, 4 (img_19, img_54, img_88).
    sw, sh = prs.slide_width, prs.slide_height
    examples_dir = os.path.join(METRIQUES_DIR, model["folder"], "output_real_examples", folder)
    all_images = sorted(glob.glob(os.path.join(examples_dir, "*.png")))
    if not all_images:
        return
    # Selection : 1ere, milieu, derniere
    if len(all_images) >= 3:
        picked = [all_images[0], all_images[len(all_images) // 2], all_images[-1]]
    else:
        picked = all_images

    # --- Slide 1 : 2 premieres images ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    title1 = f"{model['title']} - {FOLDER_LABELS[folder]} (1/2)"
    add_slide_title(slide1, title1, sw)

    margin_x = Inches(0.3)
    top_start = Inches(1.0)
    available_h = sh - top_start - Inches(0.2)
    img_h = available_h // 2
    img_w = sw - 2 * margin_x

    for i, img_path in enumerate(picked[:2]):
        y = top_start + (img_h * i)
        _add_image_keep_ratio(slide1, img_path, margin_x, y, img_w, img_h - Inches(0.05))

    if len(picked) <= 2:
        return

    # --- Slide 2 : 3eme image, plus grande ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    title2 = f"{model['title']} - {FOLDER_LABELS[folder]} (2/2)"
    add_slide_title(slide2, title2, sw)

    img_h_large = sh - top_start - Inches(0.4)
    _add_image_keep_ratio(slide2, picked[2], margin_x, top_start, img_w, img_h_large)


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    add_slide_title(slide, "Synthese comparative", sw)

    # Tableau
    rows, cols = 4, 5
    table_left = Inches(0.5)
    table_top = Inches(1.2)
    table_w = sw - Inches(1)
    table_h = Inches(2.5)
    table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_w, table_h)
    table = table_shape.table

    headers = ["Modele", "Best epoch", "Val Dice", "Val IoU", "Particularite"]
    rows_data = [
        ["NPZ Full-Res 1024", "28/30", "0.7888", "0.6717", "Image entiere redimensionnee"],
        ["NPZ Patches 256", "23/30", "0.6567", "0.5305", "Resolution native, fenetres glissantes"],
        ["Grid Major", "29/30", "0.8491", "0.7532", "Mask continu (lignes) au lieu de points"],
    ]

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        para = cell.text_frame.paragraphs[0]
        for run in para.runs:
            run.font.bold = True; run.font.size = Pt(13); run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY

    for i, row in enumerate(rows_data, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12); run.font.color.rgb = GREY

    # Conclusion textuelle
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), sw - Inches(1), Inches(2.8))
    tf = tb.text_frame; tf.word_wrap = True

    lines = [
        ("Observations", True, 18, NAVY),
        ("", False, 4, NAVY),
        ("•  Le modele Grid Major obtient le meilleur Val Dice (0.85), mais la metrique n'est pas", False, 13, GREY),
        ("   directement comparable car il predit des lignes continues, pas des points discrets.", False, 13, GREY),
        ("", False, 4, NAVY),
        ("•  Le modele Patch-Based affiche un Dice plus bas (0.66) mais opere a la resolution", False, 13, GREY),
        ("   native, ce qui rend la metrique plus exigeante (pas de redimensionnement).", False, 13, GREY),
        ("", False, 4, NAVY),
        ("•  L'evaluation visuelle sur output_real (PM Cardio) reste l'indicateur principal", False, 13, GREY),
        ("   de generalisation au reel - voir les slides exemples.", False, 13, GREY),
    ]
    for i, (text, bold, size, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("[*] Slide titre...")
    add_title_slide(prs)

    print("[*] Slide synthese...")
    add_summary_slide(prs)

    for model in MODELS:
        print(f"[*] Section : {model['title']}")
        add_section_header(prs, model["title"], model["subtitle"])
        add_metrics_slide(prs, model)
        for folder in FOLDERS:
            add_examples_slides(prs, model, folder)

    prs.save(OUTPUT_PPTX)
    print(f"\n[OK] PPTX genere : {OUTPUT_PPTX}")
    size_kb = os.path.getsize(OUTPUT_PPTX) / 1024
    print(f"     Taille : {size_kb:.0f} KB")


if __name__ == "__main__":
    main()
